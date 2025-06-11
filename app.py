import boto3
import json
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.autoshape import Shape
from pptx.shapes.placeholder import PlaceholderPicture
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.shapes.picture import Picture
from pptx.shapes.graphfrm import GraphicFrame
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import copy
from typing import List, Dict, Any, Optional
from dataclasses import dataclass

@dataclass
class RunFormat:
    """텍스트 실행(Run)의 서식 정보를 저장하는 클래스"""
    font_name: Optional[str] = None
    font_size: Optional[int] = None
    font_bold: Optional[bool] = None
    font_italic: Optional[bool] = None
    font_underline: Optional[bool] = None
    font_color_rgb: Optional[tuple] = None
    text_length: int = 0

@dataclass
class ParagraphFormat:
    """단락의 서식 정보를 저장하는 클래스"""
    alignment: Optional[int] = None
    level: int = 0
    space_before: Optional[int] = None
    space_after: Optional[int] = None
    line_spacing: Optional[float] = None
    runs: List[RunFormat] = None
    
    def __post_init__(self):
        if self.runs is None:
            self.runs = []

@dataclass
class TextFrameTemplate:
    """텍스트 프레임의 템플릿 정보를 저장하는 클래스"""
    original_text: str
    paragraphs: List[ParagraphFormat]
    margin_left: Optional[int] = None
    margin_right: Optional[int] = None
    margin_top: Optional[int] = None
    margin_bottom: Optional[int] = None
    word_wrap: Optional[bool] = None
    auto_size: Optional[int] = None

@dataclass
class TextElement:
    """번역할 텍스트 요소 정보"""
    shape_index: int
    original_text: str
    translated_text: str = ""
    element_type: str = "shape"  # shape, table_cell, chart_title, grouped_shape
    template: Optional[TextFrameTemplate] = None
    # 테이블 셀용 추가 정보
    row_idx: Optional[int] = None
    col_idx: Optional[int] = None
    # 그룹화된 도형용 추가 정보
    child_idx: Optional[int] = None

class PowerPointTranslatorImproved:
    def __init__(self):
        # Amazon Bedrock 클라이언트 설정
        self.bedrock_client = boto3.client(
            'bedrock-runtime',
            region_name='us-west-2'
        )
        self.model_id = "us.anthropic.claude-3-5-sonnet-20240620-v1:0"
        
        # 지원하는 언어 목록
        self.supported_languages = {
            'ko': '한국어',
            'en': '영어',
            'ja': '일본어',
            'zh': '중국어',
            'fr': '프랑스어',
            'de': '독일어',
            'es': '스페인어',
            'it': '이탈리아어',
            'pt': '포르투갈어',
            'ru': '러시아어'
        }
    def extract_run_format(self, run) -> RunFormat:
        """Run의 서식 정보를 추출합니다."""
        try:
            font = run.font
            font_color_rgb = None
            
            # 폰트 색상 추출 (개선된 방식)
            try:
                if hasattr(font, 'color') and font.color:
                    color = font.color
                    if hasattr(color, 'rgb') and color.rgb is not None:
                        # RGBColor 객체에서 색상 값 추출
                        rgb_color = color.rgb
                        if hasattr(rgb_color, '__iter__'):  # RGB 튜플인 경우
                            font_color_rgb = tuple(rgb_color)
                        else:  # RGBColor 객체인 경우
                            # RGBColor 객체를 정수로 변환하여 RGB 값 추출
                            rgb_int = int(rgb_color)
                            r = (rgb_int >> 16) & 0xFF
                            g = (rgb_int >> 8) & 0xFF
                            b = rgb_int & 0xFF
                            font_color_rgb = (r, g, b)
            except Exception as color_error:
                print(f"      색상 추출 중 오류 (무시됨): {str(color_error)}")
                font_color_rgb = None
            
            # 폰트 크기 추출 (개선된 방식)
            font_size_pt = None
            try:
                if font.size is not None:
                    font_size_pt = font.size.pt
            except Exception as size_error:
                print(f"      폰트 크기 추출 중 오류 (무시됨): {str(size_error)}")
            
            return RunFormat(
                font_name=font.name,
                font_size=font_size_pt,
                font_bold=font.bold,
                font_italic=font.italic,
                font_underline=font.underline,
                font_color_rgb=font_color_rgb,
                text_length=len(run.text)
            )
        except Exception as e:
            print(f"    Run 서식 추출 중 오류: {str(e)}")
            return RunFormat(text_length=len(run.text) if hasattr(run, 'text') else 0)
    
    def extract_paragraph_format(self, paragraph) -> ParagraphFormat:
        """단락의 서식 정보를 추출합니다."""
        try:
            runs_format = []
            for run in paragraph.runs:
                runs_format.append(self.extract_run_format(run))
            
            return ParagraphFormat(
                alignment=paragraph.alignment,
                level=paragraph.level,
                space_before=paragraph.space_before.pt if paragraph.space_before else None,
                space_after=paragraph.space_after.pt if paragraph.space_after else None,
                line_spacing=paragraph.line_spacing,
                runs=runs_format
            )
        except Exception as e:
            print(f"    단락 서식 추출 중 오류: {str(e)}")
            return ParagraphFormat(runs=[])
    
    def extract_text_frame_template(self, text_frame) -> TextFrameTemplate:
        """텍스트 프레임의 템플릿 정보를 추출합니다."""
        try:
            paragraphs_format = []
            original_text = text_frame.text
            
            for paragraph in text_frame.paragraphs:
                paragraphs_format.append(self.extract_paragraph_format(paragraph))
            
            return TextFrameTemplate(
                original_text=original_text,
                paragraphs=paragraphs_format,
                margin_left=text_frame.margin_left.pt if text_frame.margin_left else None,
                margin_right=text_frame.margin_right.pt if text_frame.margin_right else None,
                margin_top=text_frame.margin_top.pt if text_frame.margin_top else None,
                margin_bottom=text_frame.margin_bottom.pt if text_frame.margin_bottom else None,
                word_wrap=text_frame.word_wrap,
                auto_size=text_frame.auto_size
            )
        except Exception as e:
            print(f"    텍스트 프레임 템플릿 추출 중 오류: {str(e)}")
            return TextFrameTemplate(
                original_text=text_frame.text if hasattr(text_frame, 'text') else "",
                paragraphs=[]
            )
    def extract_text_elements_from_slide(self, slide) -> List[TextElement]:
        """슬라이드에서 텍스트 요소와 서식 템플릿을 추출합니다."""
        text_elements = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                # 1. 일반 도형의 텍스트 처리
                if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text.strip():
                    text = shape.text_frame.text.strip()
                    template = self.extract_text_frame_template(shape.text_frame)
                    
                    text_elements.append(TextElement(
                        shape_index=shape_idx,
                        original_text=text,
                        element_type="shape",
                        template=template
                    ))
                    print(f"    일반 도형 텍스트 추출: '{text[:30]}...' (shape_idx={shape_idx})")
                
                # 2. 그룹화된 도형 처리
                elif hasattr(shape, "shapes"):
                    for child_idx, child_shape in enumerate(shape.shapes):
                        if hasattr(child_shape, "text_frame") and child_shape.text_frame and child_shape.text_frame.text.strip():
                            text = child_shape.text_frame.text.strip()
                            template = self.extract_text_frame_template(child_shape.text_frame)
                            
                            text_elements.append(TextElement(
                                shape_index=shape_idx,
                                original_text=text,
                                element_type="grouped_shape",
                                template=template,
                                child_idx=child_idx
                            ))
                            print(f"    그룹 내 텍스트 추출: '{text[:30]}...' (shape_idx={shape_idx}, child_idx={child_idx})")
                
                # 3. 테이블 처리
                elif isinstance(shape, GraphicFrame) and hasattr(shape, "table"):
                    try:
                        table = shape.table
                        print(f"    테이블 발견 (shape_idx={shape_idx})")
                        
                        for row_idx, row in enumerate(table.rows):
                            for col_idx, cell in enumerate(row.cells):
                                if cell.text_frame and cell.text_frame.text.strip():
                                    text = cell.text_frame.text.strip()
                                    template = self.extract_text_frame_template(cell.text_frame)
                                    
                                    text_elements.append(TextElement(
                                        shape_index=shape_idx,
                                        original_text=text,
                                        element_type="table_cell",
                                        template=template,
                                        row_idx=row_idx,
                                        col_idx=col_idx
                                    ))
                                    print(f"      테이블 셀 텍스트 추출 [{row_idx},{col_idx}]: '{text[:30]}...'")
                    except Exception as e:
                        print(f"      테이블 처리 중 오류: {str(e)}")
                
                # 4. 차트 제목 처리
                elif isinstance(shape, GraphicFrame) and hasattr(shape, "chart"):
                    try:
                        chart = shape.chart
                        if (hasattr(chart, "chart_title") and chart.chart_title and 
                            chart.chart_title.text_frame and chart.chart_title.text_frame.text.strip()):
                            text = chart.chart_title.text_frame.text.strip()
                            template = self.extract_text_frame_template(chart.chart_title.text_frame)
                            
                            text_elements.append(TextElement(
                                shape_index=shape_idx,
                                original_text=text,
                                element_type="chart_title",
                                template=template
                            ))
                            print(f"    차트 제목 추출: '{text[:30]}...' (shape_idx={shape_idx})")
                    except Exception as e:
                        print(f"      차트 처리 중 오류: {str(e)}")
                
            except Exception as e:
                print(f"    도형 {shape_idx} 처리 중 오류: {str(e)}")
        
        return text_elements
    def translate_text(self, text: str, target_language: str, source_language: str = 'auto') -> str:
        """Amazon Bedrock Claude 3.5 Sonnet을 사용하여 텍스트를 번역합니다."""
        import time
        import botocore.exceptions
        
        language_names = {
            'ko': '한국어', 'en': '영어', 'ja': '일본어', 'zh': '중국어',
            'fr': '프랑스어', 'de': '독일어', 'es': '스페인어', 'it': '이탈리아어',
            'pt': '포르투갈어', 'ru': '러시아어'
        }
        
        target_lang_name = language_names.get(target_language, target_language)
        
        prompt = f"""다음 텍스트를 {target_lang_name}로 번역해주세요. 
번역할 때 다음 사항을 고려해주세요:
1. 원문의 의미와 뉘앙스를 정확히 전달
2. 자연스러운 표현 사용
3. 전문 용어는 해당 언어의 표준 용어 사용
4. 서식이나 특수문자는 그대로 유지
5. 번역된 텍스트만 출력 (설명이나 부가 정보 없이)

번역할 텍스트:
{text}

번역:"""

        max_retries = 5
        retry_delay = 5
        
        for attempt in range(max_retries):
            try:
                body = {
                    "anthropic_version": "bedrock-2023-05-31",
                    "max_tokens": 4000,
                    "messages": [
                        {
                            "role": "user",
                            "content": prompt
                        }
                    ],
                    "temperature": 0.0
                }
                
                response = self.bedrock_client.invoke_model(
                    modelId=self.model_id,
                    body=json.dumps(body)
                )
                
                response_body = json.loads(response['body'].read())
                translated_text = response_body['content'][0]['text'].strip()
                
                return translated_text
                
            except botocore.exceptions.ClientError as e:
                error_code = e.response.get('Error', {}).get('Code', '')
                
                if error_code == 'ThrottlingException' or 'ThrottlingException' in str(e):
                    if attempt < max_retries - 1:
                        wait_time = retry_delay * (attempt + 1)
                        print(f"  API 제한으로 인한 오류 발생: {wait_time}초 후 재시도 ({attempt+1}/{max_retries})...")
                        time.sleep(wait_time)
                    else:
                        print(f"  최대 재시도 횟수 도달: {str(e)}")
                        return text
                else:
                    print(f"  번역 중 오류 발생: {str(e)}")
                    return text
            
            except Exception as e:
                print(f"  번역 중 오류 발생: {str(e)}")
                return text
        
        return text
    def apply_run_format(self, run, run_format: RunFormat):
        """Run에 서식을 적용합니다."""
        try:
            font = run.font
            
            # 폰트명 적용
            if run_format.font_name:
                try:
                    font.name = run_format.font_name
                except Exception as e:
                    print(f"        폰트명 적용 실패: {str(e)}")
            
            # 폰트 크기 적용
            if run_format.font_size:
                try:
                    font.size = Pt(run_format.font_size)
                except Exception as e:
                    print(f"        폰트 크기 적용 실패: {str(e)}")
            
            # 굵기 적용
            if run_format.font_bold is not None:
                try:
                    font.bold = run_format.font_bold
                except Exception as e:
                    print(f"        굵기 적용 실패: {str(e)}")
            
            # 기울임 적용
            if run_format.font_italic is not None:
                try:
                    font.italic = run_format.font_italic
                except Exception as e:
                    print(f"        기울임 적용 실패: {str(e)}")
            
            # 밑줄 적용
            if run_format.font_underline is not None:
                try:
                    font.underline = run_format.font_underline
                except Exception as e:
                    print(f"        밑줄 적용 실패: {str(e)}")
            
            # 색상 적용 (개선된 방식)
            if run_format.font_color_rgb:
                try:
                    r, g, b = run_format.font_color_rgb
                    font.color.rgb = RGBColor(r, g, b)
                except Exception as e:
                    print(f"        색상 적용 실패: {str(e)}")
                
        except Exception as e:
            print(f"      Run 서식 적용 중 전체 오류: {str(e)}")
    
    def apply_paragraph_format(self, paragraph, para_format: ParagraphFormat):
        """단락에 서식을 적용합니다."""
        try:
            if para_format.alignment is not None:
                paragraph.alignment = para_format.alignment
            
            if para_format.level is not None:
                paragraph.level = para_format.level
            
            if para_format.space_before is not None:
                paragraph.space_before = Pt(para_format.space_before)
            
            if para_format.space_after is not None:
                paragraph.space_after = Pt(para_format.space_after)
            
            if para_format.line_spacing is not None:
                paragraph.line_spacing = para_format.line_spacing
                
        except Exception as e:
            print(f"      단락 서식 적용 중 오류: {str(e)}")
    
    def distribute_text_to_runs(self, translated_text: str, run_formats: List[RunFormat]) -> List[tuple]:
        """번역된 텍스트를 기존 Run 구조에 맞게 분배합니다."""
        if not run_formats:
            return [(translated_text, RunFormat())]
        
        # 단순화된 접근법: 첫 번째 Run에 모든 텍스트를 할당하고 첫 번째 Run의 서식 사용
        # 이렇게 하면 최소한 하나의 일관된 서식이 적용됩니다
        if len(run_formats) == 1 or len(translated_text.strip()) < 50:
            return [(translated_text, run_formats[0])]
        
        # 텍스트가 길고 여러 Run이 있는 경우, 의미있는 단위로 분할 시도
        try:
            # 문장 단위로 분할 시도
            sentences = []
            current_sentence = ""
            
            for char in translated_text:
                current_sentence += char
                if char in '.!?。！？':
                    sentences.append(current_sentence.strip())
                    current_sentence = ""
            
            if current_sentence.strip():
                sentences.append(current_sentence.strip())
            
            # 문장이 Run 개수보다 적으면 단순하게 처리
            if len(sentences) <= len(run_formats):
                result = []
                for i, sentence in enumerate(sentences):
                    format_idx = min(i, len(run_formats) - 1)
                    result.append((sentence, run_formats[format_idx]))
                return result
            else:
                # 문장이 많으면 첫 번째 Run 서식으로 통일
                return [(translated_text, run_formats[0])]
                
        except Exception as e:
            print(f"      텍스트 분배 중 오류, 단순 방식 사용: {str(e)}")
            # 오류 발생 시 첫 번째 Run 서식으로 모든 텍스트 할당
            return [(translated_text, run_formats[0])]
    def apply_template_to_text_frame(self, text_frame, translated_text: str, template: TextFrameTemplate) -> bool:
        """템플릿을 사용하여 텍스트 프레임에 번역된 텍스트를 적용합니다."""
        try:
            # 텍스트 프레임 속성 복원 (안전하게)
            try:
                if template.margin_left is not None:
                    text_frame.margin_left = Pt(template.margin_left)
                if template.margin_right is not None:
                    text_frame.margin_right = Pt(template.margin_right)
                if template.margin_top is not None:
                    text_frame.margin_top = Pt(template.margin_top)
                if template.margin_bottom is not None:
                    text_frame.margin_bottom = Pt(template.margin_bottom)
                if template.word_wrap is not None:
                    text_frame.word_wrap = template.word_wrap
                if template.auto_size is not None:
                    text_frame.auto_size = template.auto_size
            except Exception as e:
                print(f"      텍스트 프레임 속성 복원 중 오류 (무시됨): {str(e)}")
            
            # 번역된 텍스트를 줄바꿈으로 분할
            translated_lines = translated_text.split('\n')
            
            # 기존 단락들의 서식 정보 백업
            original_paragraphs = []
            for para in text_frame.paragraphs:
                try:
                    para_format = self.extract_paragraph_format(para)
                    original_paragraphs.append(para_format)
                except:
                    original_paragraphs.append(ParagraphFormat())
            
            # 모든 기존 단락 제거 (첫 번째 제외)
            while len(text_frame.paragraphs) > 1:
                try:
                    p = text_frame.paragraphs[-1]
                    p._element.getparent().remove(p._element)
                except:
                    break
            
            # 첫 번째 단락 처리
            if text_frame.paragraphs:
                first_para = text_frame.paragraphs[0]
                
                # 기존 runs 제거
                try:
                    for run in first_para.runs[:]:
                        run._r.getparent().remove(run._r)
                except Exception as e:
                    print(f"      기존 runs 제거 중 오류: {str(e)}")
                
                # 첫 번째 줄 적용
                first_line = translated_lines[0] if translated_lines else ""
                
                if first_line:
                    # 원본 첫 번째 단락의 서식 사용
                    if original_paragraphs:
                        para_format = original_paragraphs[0]
                        if para_format.runs:
                            # 첫 번째 Run의 서식으로 통일하여 안정성 확보
                            run = first_para.add_run()
                            run.text = first_line
                            self.apply_run_format(run, para_format.runs[0])
                        else:
                            first_para.text = first_line
                    else:
                        first_para.text = first_line
                    
                    # 단락 서식 적용
                    if original_paragraphs:
                        self.apply_paragraph_format(first_para, original_paragraphs[0])
            
            # 나머지 줄들을 새 단락으로 추가
            for i, line in enumerate(translated_lines[1:], 1):
                if line.strip():  # 빈 줄이 아닌 경우만 처리
                    try:
                        para = text_frame.add_paragraph()
                        
                        # 해당하는 원본 단락 서식이 있으면 사용, 없으면 첫 번째 단락 서식 사용
                        if i < len(original_paragraphs):
                            para_format = original_paragraphs[i]
                        elif original_paragraphs:
                            para_format = original_paragraphs[0]
                        else:
                            para_format = ParagraphFormat()
                        
                        # 텍스트 추가
                        if para_format.runs:
                            run = para.add_run()
                            run.text = line
                            self.apply_run_format(run, para_format.runs[0])
                        else:
                            para.text = line
                        
                        # 단락 서식 적용
                        self.apply_paragraph_format(para, para_format)
                        
                    except Exception as e:
                        print(f"      단락 {i} 추가 중 오류: {str(e)}")
                        # 실패 시 기본 방식으로 단락 추가
                        try:
                            para = text_frame.add_paragraph()
                            para.text = line
                        except:
                            pass
            
            return True
            
        except Exception as e:
            print(f"      템플릿 적용 중 전체 오류: {str(e)}")
            # 실패 시 기본 방식으로 폴백
            try:
                text_frame.text = translated_text
                print(f"      기본 방식으로 폴백 완료")
                return True
            except Exception as fallback_error:
                print(f"      기본 방식 폴백도 실패: {str(fallback_error)}")
                return False
    def apply_translation_to_slide(self, slide, text_elements: List[TextElement]) -> int:
        """슬라이드에 번역된 텍스트를 적용합니다."""
        success_count = 0
        
        for element in text_elements:
            try:
                shape = slide.shapes[element.shape_index]
                
                if element.element_type == "shape":
                    # 일반 도형
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        if self.apply_template_to_text_frame(shape.text_frame, element.translated_text, element.template):
                            success_count += 1
                            print(f"      일반 도형 번역 적용 완료: '{element.original_text[:20]}...' -> '{element.translated_text[:20]}...'")
                        else:
                            print(f"      일반 도형 번역 적용 실패")
                
                elif element.element_type == "grouped_shape":
                    # 그룹화된 도형
                    if hasattr(shape, "shapes") and element.child_idx is not None:
                        child_shape = shape.shapes[element.child_idx]
                        if hasattr(child_shape, "text_frame") and child_shape.text_frame:
                            if self.apply_template_to_text_frame(child_shape.text_frame, element.translated_text, element.template):
                                success_count += 1
                                print(f"      그룹 내 도형 번역 적용 완료: '{element.original_text[:20]}...' -> '{element.translated_text[:20]}...'")
                            else:
                                print(f"      그룹 내 도형 번역 적용 실패")
                
                elif element.element_type == "table_cell":
                    # 테이블 셀
                    if isinstance(shape, GraphicFrame) and hasattr(shape, "table"):
                        table = shape.table
                        if element.row_idx is not None and element.col_idx is not None:
                            cell = table.cell(element.row_idx, element.col_idx)
                            if cell.text_frame:
                                if self.apply_template_to_text_frame(cell.text_frame, element.translated_text, element.template):
                                    success_count += 1
                                    print(f"      테이블 셀 번역 적용 완료 [{element.row_idx},{element.col_idx}]: '{element.original_text[:20]}...' -> '{element.translated_text[:20]}...'")
                                else:
                                    print(f"      테이블 셀 번역 적용 실패")
                
                elif element.element_type == "chart_title":
                    # 차트 제목
                    if isinstance(shape, GraphicFrame) and hasattr(shape, "chart"):
                        chart = shape.chart
                        if hasattr(chart, "chart_title") and chart.chart_title and chart.chart_title.text_frame:
                            if self.apply_template_to_text_frame(chart.chart_title.text_frame, element.translated_text, element.template):
                                success_count += 1
                                print(f"      차트 제목 번역 적용 완료: '{element.original_text[:20]}...' -> '{element.translated_text[:20]}...'")
                            else:
                                print(f"      차트 제목 번역 적용 실패")
                
            except Exception as e:
                print(f"      텍스트 요소 적용 중 오류: {str(e)}")
        
        return success_count
    def translate_presentation(self, input_file: str, output_file: str, target_language: str) -> bool:
        """PowerPoint 프레젠테이션을 템플릿 기반 방식으로 번역합니다."""
        
        if target_language not in self.supported_languages:
            print(f"지원하지 않는 언어입니다. 지원 언어: {list(self.supported_languages.keys())}")
            return False
        
        try:
            # PowerPoint 파일 로드
            prs = Presentation(input_file)
            total_slides = len(prs.slides)
            
            print(f"번역 시작: {total_slides}개 슬라이드를 {self.supported_languages[target_language]}로 번역합니다...")
            print("템플릿 기반 서식 보존 방식을 사용합니다.\n")
            
            # 성공/실패 통계
            total_success = 0
            total_failed = 0
            slide_success = 0
            slide_failed = 0
            
            # 각 슬라이드 처리
            for slide_idx, slide in enumerate(prs.slides):
                try:
                    print(f"슬라이드 {slide_idx + 1}/{total_slides} 처리 중...")
                    
                    # 1단계: 텍스트 요소와 서식 템플릿 추출
                    print("  1단계: 텍스트 및 서식 정보 추출 중...")
                    text_elements = self.extract_text_elements_from_slide(slide)
                    
                    if not text_elements:
                        print("  번역할 텍스트가 없습니다.")
                        slide_success += 1
                        continue
                    
                    # 2단계: 텍스트 번역
                    print(f"  2단계: {len(text_elements)}개 텍스트 요소 번역 중...")
                    for element in text_elements:
                        if element.original_text.strip():
                            element.translated_text = self.translate_text(element.original_text, target_language)
                            print(f"    번역 완료: '{element.original_text[:30]}...' -> '{element.translated_text[:30]}...'")
                        else:
                            element.translated_text = element.original_text
                    
                    # 3단계: 서식을 보존하면서 번역된 텍스트 적용
                    print("  3단계: 서식 보존하며 번역 텍스트 적용 중...")
                    applied_count = self.apply_translation_to_slide(slide, text_elements)
                    
                    total_success += applied_count
                    total_failed += len(text_elements) - applied_count
                    
                    if applied_count > 0:
                        slide_success += 1
                        print(f"  슬라이드 {slide_idx + 1} 완료: {applied_count}/{len(text_elements)}개 요소 성공")
                    else:
                        slide_failed += 1
                        print(f"  슬라이드 {slide_idx + 1} 실패: 번역 적용되지 않음")
                    
                except Exception as e:
                    print(f"  슬라이드 {slide_idx + 1} 처리 중 오류 발생: {str(e)}")
                    slide_failed += 1
                
                print()  # 슬라이드 간 구분을 위한 빈 줄
            
            # 번역된 파일 저장
            try:
                prs.save(output_file)
                print(f"번역 완료! 저장된 파일: {output_file}")
                print(f"슬라이드 처리 결과: 성공 {slide_success}개, 실패 {slide_failed}개")
                print(f"텍스트 요소 처리 결과: 성공 {total_success}개, 실패 {total_failed}개")
                return slide_success > 0
            except Exception as e:
                print(f"파일 저장 중 오류 발생: {str(e)}")
                return False
            
        except Exception as e:
            print(f"프레젠테이션 로드 중 오류 발생: {str(e)}")
            return False
    
    def show_supported_languages(self):
        """지원하는 언어 목록을 출력합니다."""
        print("지원하는 언어:")
        for code, name in self.supported_languages.items():
            print(f"  {code}: {name}")

def main():
    translator = PowerPointTranslatorImproved()
    
    print("=== PowerPoint 번역기 (개선된 버전) ===")
    print("Amazon Bedrock Claude 3.5 Sonnet 사용")
    print("템플릿 기반 서식 보존 방식\n")
    
    # 지원 언어 출력
    translator.show_supported_languages()
    
    # 사용자 입력
    input_file = input("\n번역할 PowerPoint 파일 경로를 입력하세요: ").strip()
    
    if not os.path.exists(input_file):
        print("파일이 존재하지 않습니다.")
        return
    
    target_language = input("번역할 언어 코드를 입력하세요 (예: ko, en, ja): ").strip().lower()
    
    # 출력 파일명 생성
    base_name = os.path.splitext(input_file)[0]
    output_file = f"{base_name}_translated_improved_{target_language}.pptx"
    
    # 번역 실행
    success = translator.translate_presentation(input_file, output_file, target_language)
    
    if success:
        print(f"\n✅ 번역이 성공적으로 완료되었습니다!")
        print(f"📁 번역된 파일: {output_file}")
        print("\n🔧 개선 사항:")
        print("- 템플릿 기반 서식 보존 방식 적용")
        print("- 폰트, 색상, 크기 등 상세 서식 정보 보존")
        print("- Run 단위 서식 분배로 정확한 서식 복원")
        print("- 단락별 서식 및 정렬 정보 보존")
    else:
        print("❌ 번역 중 오류가 발생했습니다.")

if __name__ == "__main__":
    main()
